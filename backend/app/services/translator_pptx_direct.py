#!/usr/bin/env python3
"""
translator_pptx_direct.py

并行优化版 PPTX 翻译器（合并 run、无调试标记）

Usage:
    python translator_pptx_direct.py input.pptx output.pptx src_lang tgt_lang [engine]

依赖:
    pip install python-pptx
    并确保项目内有 translate_batch(texts, src, tgt, engine=...) 可调用
"""
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
import sys
import math
from concurrent.futures import ThreadPoolExecutor, as_completed
from typing import List, Tuple, Dict, Any
import traceback

# 调整下面导入为你项目中 utils_translator 的位置
try:
    # prefer package import if running inside app module
    from app.services.utils_translator import translate_batch
    from app.services.terminology_service import (
        get_terminology_options,
        preprocess_texts,
        preprocess_texts_with_categories,
        postprocess_texts,
    )
    from app.database import SessionLocal
except Exception:
    # fallback: assume utils_translator.py is in same folder
    try:
        from utils_translator import translate_batch
    except Exception:
        raise ImportError(
            "无法导入 translate_batch。请确保 app.services.utils_translator.translate_batch 可用，"
            "或者把 utils_translator.py 放到同级目录。"
        )


def is_translatable(text: str) -> bool:
    """判断是否需要翻译：空/纯数字/纯符号/过短文本跳过"""
    if text is None:
        return False
    s = text.strip()
    if not s:
        return False
    if s.isdigit():
        return False
    # 如果都是标点/符号（没有字母或汉字等），跳过
    if all(not ch.isalnum() for ch in s):
        return False
    # 过滤掉过短的文本（少于2个字符），避免翻译失败
    if len(s) < 2:
        return False
    # 过滤掉只包含单个字符的文本
    if len(s) == 1 and not s.isalnum():
        return False
    return True


def chunk_list(lst: List[Any], n_chunks: int) -> List[List[Any]]:
    """把 list 平均切分成 n_chunks 片（尽量均匀）"""
    if n_chunks <= 0:
        return [lst]
    k = max(1, math.ceil(len(lst) / n_chunks))
    return [lst[i:i + k] for i in range(0, len(lst), k)]


def batch_translate_parallel(texts: List[str], src: str, tgt: str, engine: str = "deepseek", max_workers: int = 4, **options) -> List[str]:
    """
    将 texts 切块后并行调用 translate_batch，对每个块返回翻译结果（保持 order）
    如果某个块抛异常，会在主线程中尝试二分重试（递归）。
    """
    if not texts:
        return []

    # small optimization: if only one text, call directly and unwrap result
    if len(texts) == 1:
        _res = translate_batch(texts, src, tgt, engine=engine, **options)
        if isinstance(_res, tuple) and len(_res) >= 1:
            _res = _res[0]
        # ensure list[str]
        if isinstance(_res, list):
            return [str(_res[0])] if _res else [texts[0]]
        return [str(_res)]

    # choose number of chunks by length and max_workers
    n_chunks = min(max_workers, max(1, len(texts) // 8))  # 每 chunk 最少包含 ~8 个句子以减少请求数
    chunks = chunk_list(texts, n_chunks)

    results = [None] * len(texts)
    index_offsets = []
    offset = 0
    for c in chunks:
        index_offsets.append(offset)
        offset += len(c)

    # ThreadPoolExecutor calling translate_batch for each chunk
    with ThreadPoolExecutor(max_workers=max_workers) as ex:
        future_to_chunk_idx = {}
        for idx, chunk in enumerate(chunks):
            future = ex.submit(translate_batch, chunk, src, tgt, engine, **options)
            future_to_chunk_idx[future] = idx

        for fut in as_completed(future_to_chunk_idx):
            idx = future_to_chunk_idx[fut]
            offset = index_offsets[idx]
            try:
                translated_chunk = fut.result()
                # unwrap tuple (texts, tokens) if necessary
                if isinstance(translated_chunk, tuple) and len(translated_chunk) >= 1:
                    translated_chunk = translated_chunk[0]
                # safety: ensure returned length matches
                if not isinstance(translated_chunk, list) or len(translated_chunk) != len(chunks[idx]):
                    # fallback to single-item translate for that chunk
                    translated_chunk = []
                    for item in chunks[idx]:
                        _single = translate_batch([item], src, tgt, engine, **options)
                        if isinstance(_single, tuple) and len(_single) >= 1:
                            _single = _single[0]
                        if isinstance(_single, list) and len(_single) >= 1:
                            translated_chunk.append(str(_single[0]))
                        else:
                            translated_chunk.append(item)
                for i, t in enumerate(translated_chunk):
                    results[offset + i] = str(t) if t is not None else texts[offset + i]
            except Exception as e:
                # 异常：对该 chunk 进行二分重试
                try:
                    left = chunks[idx][:len(chunks[idx]) // 2]
                    right = chunks[idx][len(chunks[idx]) // 2:]
                    left_tr = batch_translate_parallel(left, src, tgt, engine, max_workers=max_workers)
                    right_tr = batch_translate_parallel(right, src, tgt, engine, max_workers=max_workers)
                    translated_chunk = left_tr + right_tr
                    for i, t in enumerate(translated_chunk):
                        results[offset + i] = t
                except Exception:
                    # 最后退化：逐条调用 translate_batch
                    for i, item in enumerate(chunks[idx]):
                        single_tr = translate_batch([item], src, tgt, engine, **options)
                        if isinstance(single_tr, tuple) and len(single_tr) >= 1:
                            single_tr = single_tr[0]
                        if isinstance(single_tr, list) and len(single_tr) >= 1:
                            results[offset + i] = str(single_tr[0])
                        else:
                            results[offset + i] = item

    # final safety: replace any None by original texts
    for i in range(len(results)):
        if results[i] is None:
            results[i] = texts[i]
    return results


def collect_text_items(prs: Presentation) -> Tuple[List[Tuple[str, int, int, int]], List[str]]:
    """
    遍历 PPTX 提取可翻译文本段（合并 run 后的段落）
    返回:
      items: list of (part_key, slide_idx, shape_idx, para_idx)
             part_key 用于在后续写回时定位（这里我们用 tuple 转为 str）
      texts: list of paragraph texts (same order)
    说明：part_key 只是个标识，写回时采用相同的遍历顺序来写回（保证一一对应）
    """
    items = []
    texts = []

    def iter_shapes_recursive(container):
        for shp in container.shapes:
            yield shp
            # 递归处理组合形状
            try:
                if shp.shape_type == MSO_SHAPE_TYPE.GROUP and hasattr(shp, 'shapes'):
                    for inner in iter_shapes_recursive(shp):
                        yield inner
            except Exception:
                pass

    for s_idx, slide in enumerate(prs.slides):
        for sh_idx, shape in enumerate(list(iter_shapes_recursive(slide))):
            # text frame (普通文本框/占位/标题)
            if shape.has_text_frame:
                tf = shape.text_frame
                for p_idx, para in enumerate(tf.paragraphs):
                    merged = "".join(run.text or "" for run in para.runs)
                    if is_translatable(merged):
                        items.append((s_idx, sh_idx, p_idx, "text_frame"))
                        texts.append(merged)
            # table cell
            elif shape.has_table:
                table = shape.table
                for r in range(len(table.rows)):
                    for c in range(len(table.columns)):
                        cell = table.cell(r, c)
                        # treat each paragraph in cell
                        for p_idx, para in enumerate(cell.text_frame.paragraphs):
                            merged = "".join(run.text or "" for run in para.runs)
                            if is_translatable(merged):
                                # encode cell position in item
                                items.append((s_idx, sh_idx, (r, c, p_idx), "table_cell"))
                                texts.append(merged)

        # slide notes
        if slide.has_notes_slide:
            notes = slide.notes_slide
            for sh_idx, nshape in enumerate(notes.shapes):
                if nshape.has_text_frame:
                    for p_idx, para in enumerate(nshape.text_frame.paragraphs):
                        merged = "".join(run.text or "" for run in para.runs)
                        if is_translatable(merged):
                            items.append((s_idx, sh_idx, p_idx, "notes"))
                            texts.append(merged)

    return items, texts


def write_translations_back(prs: Presentation, items: List[Tuple], original_texts: List[str], translated_map: Dict[str, str]) -> None:
    """
    根据 items（与 texts 同序）把 translations 写回到 prs 中。
    写回策略：
      - 针对每个段落，先清除该段落所有 run，然后添加一个 run 放入译文
      - 保留第一个 run 的样式（若存在）作为大致样式
      - 对表格单元格与 notes 做相同处理
    """
    for idx, item in enumerate(items):
        s_idx, sh_idx, p_info, kind = item
        translated = translated_map.get(original_texts[idx], original_texts[idx])

        slide = prs.slides[s_idx]
        # bounds checking
        if sh_idx >= len(slide.shapes):
            continue
        shape = slide.shapes[sh_idx]

        if kind == "text_frame":
            if not shape.has_text_frame:
                continue
            para_idx = p_info
            tf = shape.text_frame
            if para_idx >= len(tf.paragraphs):
                continue
            para = tf.paragraphs[para_idx]
            # capture style from first run if exists
            template_style = {}
            if para.runs and len(para.runs) > 0:
                r0 = para.runs[0]
                template_style = {
                    "font_name": r0.font.name,
                    "size": r0.font.size,
                    "bold": r0.font.bold,
                    "italic": r0.font.italic,
                    "underline": r0.font.underline,
                    "color": getattr(r0.font.color, 'rgb', None)
                }
            # clear paragraph's runs by setting text to ''
            # paragraph has no clear() API, but we can replace runs by setting text of first run and removing others
            if para.runs:
                # set first run to translated text
                para.runs[0].text = translated
                # apply template style if present
                r = para.runs[0]
                apply_run_style_from_dict(r, template_style)
                # remove extra runs
                for _ in range(len(para.runs) - 1):
                    # remove last run
                    try:
                        last = para.runs[-1]
                        # there's no direct remove; set text to '' to effectively clear it
                        last.text = ""
                    except Exception:
                        pass
            else:
                # no run exists? add one
                run = para.add_run()
                run.text = translated

        elif kind == "table_cell":
            # p_info = (row, col, para_idx)
            r_idx, c_idx, para_idx = p_info
            if not shape.has_table:
                continue
            table = shape.table
            if r_idx >= len(table.rows) or c_idx >= len(table.columns):
                continue
            cell = table.cell(r_idx, c_idx)
            tf = cell.text_frame
            if para_idx >= len(tf.paragraphs):
                continue
            para = tf.paragraphs[para_idx]
            template_style = {}
            if para.runs and len(para.runs) > 0:
                r0 = para.runs[0]
                template_style = {
                    "font_name": r0.font.name,
                    "size": r0.font.size,
                    "bold": r0.font.bold,
                    "italic": r0.font.italic,
                    "underline": r0.font.underline,
                    "color": getattr(r0.font.color, 'rgb', None)
                }
            if para.runs:
                para.runs[0].text = translated
                apply_run_style_from_dict(para.runs[0], template_style)
                for _ in range(len(para.runs) - 1):
                    try:
                        para.runs[-1].text = ""
                    except Exception:
                        pass
            else:
                run = para.add_run()
                run.text = translated

        elif kind == "notes":
            # p_info is paragraph index in notes slide shapes; we used sh_idx to find a notes shape index earlier
            para_idx = p_info
            notes = slide.notes_slide
            # find the shape at sh_idx (may be different index if shapes changed)
            if sh_idx >= len(notes.shapes):
                # fallback: try to find any text shape and use para_idx-th paragraph
                for nsh in notes.shapes:
                    if nsh.has_text_frame:
                        tf = nsh.text_frame
                        if para_idx < len(tf.paragraphs):
                            para = tf.paragraphs[para_idx]
                            if para.runs:
                                para.runs[0].text = translated
                                apply_run_style_from_dict(para.runs[0], {})
                                for _ in range(len(para.runs) - 1):
                                    para.runs[-1].text = ""
                            else:
                                run = para.add_run()
                                run.text = translated
                            break
                continue
            nshape = notes.shapes[sh_idx]
            if not nshape.has_text_frame:
                continue
            tf = nshape.text_frame
            if para_idx >= len(tf.paragraphs):
                continue
            para = tf.paragraphs[para_idx]
            template_style = {}
            if para.runs and len(para.runs) > 0:
                r0 = para.runs[0]
                template_style = {
                    "font_name": r0.font.name,
                    "size": r0.font.size,
                    "bold": r0.font.bold,
                    "italic": r0.font.italic,
                    "underline": r0.font.underline,
                    "color": getattr(r0.font.color, 'rgb', None)
                }
            if para.runs:
                para.runs[0].text = translated
                apply_run_style_from_dict(para.runs[0], template_style)
                for _ in range(len(para.runs) - 1):
                    try:
                        para.runs[-1].text = ""
                    except Exception:
                        pass
            else:
                run = para.add_run()
                run.text = translated


def apply_run_style_from_dict(run, style: Dict[str, Any]):
    """把 style 应用到 run.font"""
    font = run.font
    if not style:
        return
    if style.get("font_name"):
        try:
            font.name = style["font_name"]
        except Exception:
            pass
    if style.get("size"):
        try:
            font.size = style["size"]
        except Exception:
            pass
    if style.get("bold") is not None:
        font.bold = style["bold"]
    if style.get("italic") is not None:
        font.italic = style["italic"]
    if style.get("underline") is not None:
        font.underline = style["underline"]
    if style.get("color") is not None:
        try:
            # color is RGBColor
            if isinstance(style["color"], RGBColor):
                font.color.rgb = style["color"]
            else:
                # if it's tuple or hex string, try to set flexibly
                font.color.rgb = style["color"]
        except Exception:
            pass


def translate_pptx(input_path: str, output_path: str, src: str, tgt: str, engine: str = "deepseek", max_workers: int = 4, user_id: int | None = None, **kwargs):
    prs = Presentation(input_path)

    items, texts = collect_text_items(prs)
    if not texts:
        prs.save(output_path)
        return {"token_count": 0, "character_count": 0}

    total_character_count = sum(len(t) for t in texts)

    # 去重：保持 items/texts 顺序，但创建 unique list 与映射
    unique_texts = []
    seen = {}
    for t in texts:
        if t not in seen:
            seen[t] = len(unique_texts)
            unique_texts.append(t)

    # 并行批量翻译 unique_texts
    # 术语前/后处理
    try:
        db = SessionLocal()
        options = get_terminology_options(db)
        if options.get("terminology_enabled", True):
            processed_texts, mappings = preprocess_texts(db, unique_texts, src, tgt, case_sensitive=bool(options.get("case_sensitive", False)), user_id=user_id)
        else:
            processed_texts, mappings = unique_texts, [{} for _ in unique_texts]
    finally:
        try:
            db.close()
        except Exception:
            pass

    # Qwen3: 避免并行，改为顺序以降低 429 触发率
    if engine and str(engine).lower() == 'qwen3':
        # 顺序调用逐条翻译，并加入轻微节流，减少429
        translated_unique = []
        for s in processed_texts:
            _r = translate_batch([s], src, tgt, engine=engine)
            if isinstance(_r, tuple) and len(_r) >= 1:
                _r = _r[0]
            if isinstance(_r, list) and _r:
                translated_unique.append(str(_r[0]))
            else:
                translated_unique.append(s)
            try:
                import time as _t
                _t.sleep(0.05)
            except Exception:
                pass
    else:
        # 并行批量翻译，提升大文档速度
        translated_unique = batch_translate_parallel(processed_texts, src, tgt, engine=engine, max_workers=max_workers)
    total_token_count = 0

    if options.get("terminology_enabled", True):
        translated_unique = postprocess_texts(translated_unique, mappings)

    # build map original -> translated
    translated_map = {u: translated_unique[i] for i, u in enumerate(unique_texts)}

    # 写回所有段落（按原始 items 顺序）
    write_translations_back(prs, items, texts, translated_map)

    # 保存
    prs.save(output_path)
    
    return {
        "token_count": total_token_count,
        "character_count": total_character_count
    }

import pptx
from ..services.utils_translator import translate_batch
from .translator_pptx_ooxml import translate_pptx_ooxml

def is_translatable(text):
    """判断文本是否需要翻译：非空、非纯数字、非纯符号"""
    if not text or not text.strip():
        return False
    stripped = text.strip()
    if stripped.isdigit():
        return False
    if all(not ch.isalnum() for ch in stripped):
        return False
    return True

def _normalize_text(s: str) -> str:
    if not isinstance(s, str):
        return s
    # 清理控制字符标记，常见 _x000B_, _x000D_ 等；替换为换行或空
    s = s.replace("_x000B_", "\n").replace("_x000b_", "\n")
    s = s.replace("_x000D_", "\n").replace("_x000d_", "\n")
    # 也清理真实的 VT 控制符
    s = s.replace("\x0b", "\n")
    return s


def _iter_shapes_recursive(container):
    for shp in container.shapes:
        yield shp
        try:
            if shp.shape_type == MSO_SHAPE_TYPE.GROUP and hasattr(shp, 'shapes'):
                for inner in _iter_shapes_recursive(shp):
                    yield inner
        except Exception:
            pass


def translate_pptx_direct(input_path, output_path, src_lang, tgt_lang, engine="deepseek", user_id: int | None = None, **kwargs):
    # 优先使用 OOXML 层替换，仅改 a:t 文本，最大化保留样式/布局
    try:
        cat_ids = kwargs.get('category_ids')
        return translate_pptx_ooxml(input_path, output_path, src_lang, tgt_lang, engine=engine, user_id=user_id, category_ids=cat_ids)
    except Exception:
        # 出错时回退到 python-pptx 改写方案
        pass
    """Translate a PPTX file directly, preserving formatting and attempting layout adjustments."""
    from pptx import Presentation
    
    # Load the presentation
    prs = Presentation(input_path)
    total_token_count = 0
    total_character_count = 0
    
    # 第一步：收集所有需要翻译的文本和位置信息
    text_items = []  # [(slide_idx, shape_idx, paragraph_idx, text, is_table, row_idx, col_idx), ...]
    
    for slide_idx, slide in enumerate(prs.slides):
        # 递归遍历形状（支持组合）
        for shape_idx, shape in enumerate(list(_iter_shapes_recursive(slide))):
            if hasattr(shape, "text_frame") and shape.text_frame:
                for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    original_text = _normalize_text(paragraph.text)
                    if is_translatable(original_text):
                        text_items.append({
                            'type': 'shape',
                            'slide_idx': slide_idx,
                            'shape_idx': shape_idx,
                            'para_idx': para_idx,
                            'text': original_text,
                            'paragraph': paragraph
                        })
                        total_character_count += len(original_text)

            # Collect text from tables
            if shape.has_table:
                for row_idx, row in enumerate(shape.table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        for para_idx, paragraph in enumerate(cell.text_frame.paragraphs):
                            original_text = _normalize_text(paragraph.text)
                            if is_translatable(original_text):
                                text_items.append({
                                    'type': 'table',
                                    'slide_idx': slide_idx,
                                    'shape_idx': shape_idx,
                                    'row_idx': row_idx,
                                    'col_idx': col_idx,
                                    'para_idx': para_idx,
                                    'text': original_text,
                                    'paragraph': paragraph
                                })
                                total_character_count += len(original_text)
        
        # Collect text from slide notes
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame:
                for para_idx, paragraph in enumerate(notes_slide.notes_text_frame.paragraphs):
                    original_text = _normalize_text(paragraph.text)
                    if is_translatable(original_text):
                        text_items.append({
                            'type': 'notes',
                            'slide_idx': slide_idx,
                            'para_idx': para_idx,
                            'text': original_text,
                            'paragraph': paragraph
                        })
                        total_character_count += len(original_text)
    
    print(f"收集到 {len(text_items)} 个需要翻译的文本项，总字符数: {total_character_count}")
    
    if not text_items:
        print("没有需要翻译的文本，直接保存")
        prs.save(output_path)
        return {
            "translated_file_path": output_path,
            "token_count": 0,
            "character_count": 0
        }
    
    # 第二步：批量翻译所有文本
    print(f"开始批量翻译 {len(text_items)} 个文本...")
    all_texts = [item['text'] for item in text_items]
    
    try:
        # 使用批量翻译，大幅减少API调用次数
        # 术语前/后处理
        try:
            db = SessionLocal()
            options = get_terminology_options(db)
            if options.get("terminology_enabled", True):
                cat_ids = kwargs.get("category_ids", None)
                if options.get("categories_enabled", True):
                    if not cat_ids or (isinstance(cat_ids, list) and len(cat_ids) == 0):
                        processed_texts, mappings = all_texts, [{} for _ in all_texts]
                    else:
                        processed_texts, mappings = preprocess_texts_with_categories(
                            db, all_texts, src_lang, tgt_lang, cat_ids,
                            case_sensitive=bool(options.get("case_sensitive", False)), user_id=user_id
                        )
                else:
                    processed_texts, mappings = preprocess_texts(
                        db, all_texts, src_lang, tgt_lang,
                        case_sensitive=bool(options.get("case_sensitive", False)), user_id=user_id
                    )
            else:
                processed_texts, mappings = all_texts, [{} for _ in all_texts]
        finally:
            try:
                db.close()
            except Exception:
                pass

        translated_texts, batch_token_count = translate_batch(processed_texts, src_lang, tgt_lang, engine=engine)
        total_token_count = batch_token_count
        
        if len(translated_texts) != len(all_texts):
            print(f"警告：翻译结果数量不匹配，期望 {len(all_texts)}，实际 {len(translated_texts)}")
            # 如果数量不匹配，使用原文填充
            while len(translated_texts) < len(all_texts):
                translated_texts.append(all_texts[len(translated_texts)])
        
        # 术语后处理
        if options.get("terminology_enabled", True):
            translated_texts = postprocess_texts(translated_texts, mappings)

        print(f"批量翻译完成，消耗 {total_token_count} tokens")
        
    except Exception as e:
        print(f"批量翻译失败: {e}")
        # 直接返回失败，不进行逐段翻译回退
        print("翻译失败，返回错误")
        raise e
    
    # 第三步：将翻译结果写回PPTX
    print("开始将翻译结果写回PPTX...")
    
    # 添加调试信息
    print(f"调试信息：")
    print(f"  - 原始文本数量: {len(all_texts)}")
    print(f"  - 翻译结果数量: {len(translated_texts)}")
    print(f"  - 前5个原始文本: {all_texts[:5]}")
    print(f"  - 前5个翻译结果: {translated_texts[:5]}")
    
    # 统计翻译结果
    translated_count = 0
    untranslated_count = 0
    
    for i, item in enumerate(text_items):
        if i < len(translated_texts):
            translated_text = translated_texts[i]
            original_text = item['text']
            
            # 添加详细的调试日志
            if i < 5:  # 只显示前5个的详细信息
                print(f"调试 - 项目 {i}:")
                print(f"  原始文本: '{original_text}'")
                print(f"  翻译结果: '{translated_text}'")
                print(f"  类型: {type(translated_text)}")
                print(f"  是否相等: {translated_text == original_text}")
            
            # 检查翻译结果的有效性
            if translated_text is None or not isinstance(translated_text, str):
                print(f"警告: 文本 '{original_text[:50]}...' 的翻译结果无效，使用原文")
                translated_text = original_text
            elif translated_text.strip() == "":
                print(f"警告: 文本 '{original_text[:50]}...' 的翻译结果为空，使用原文")
                translated_text = original_text
            else:
                translated_text = _normalize_text(translated_text)
            
            # 检查是否真的被翻译了（不是原文）
            if translated_text != original_text:
                translated_count += 1
                if i < 5:  # 只显示前5个的详细信息
                    print(f"  ✅ 翻译成功: '{original_text[:30]}...' -> '{translated_text[:30]}...'")
            else:
                untranslated_count += 1
                if i < 5:  # 只显示前5个的详细信息
                    print(f"  ❌ 翻译失败或未变化: '{original_text[:30]}...'")
            
            # 确保translated_text是有效的字符串
            if not isinstance(translated_text, str):
                translated_text = str(translated_text) if translated_text is not None else original_text
            
            if item['type'] == 'shape':
                paragraph = item.get('paragraph') or prs.slides[item['slide_idx']].shapes[item['shape_idx']].text_frame.paragraphs[item['para_idx']]
                # 尽量保留样式：保留第一个 run，写入译文，清空其余 run 文本
                if paragraph.runs:
                    paragraph.runs[0].text = translated_text
                    # 其余 run 清空文本但不移除，保留样式
                    for r in paragraph.runs[1:]:
                        r.text = ""
                else:
                    paragraph.add_run().text = translated_text
                
            elif item['type'] == 'table':
                paragraph = item.get('paragraph')
                if paragraph is None:
                    slide = prs.slides[item['slide_idx']]
                    shape = slide.shapes[item['shape_idx']]
                    cell = shape.table.rows[item['row_idx']].cells[item['col_idx']]
                    paragraph = cell.text_frame.paragraphs[item['para_idx']]
                if paragraph.runs:
                    paragraph.runs[0].text = translated_text
                    for r in paragraph.runs[1:]:
                        r.text = ""
                else:
                    paragraph.add_run().text = translated_text
                
            elif item['type'] == 'notes':
                paragraph = item.get('paragraph') or prs.slides[item['slide_idx']].notes_slide.notes_text_frame.paragraphs[item['para_idx']]
                if paragraph.runs:
                    paragraph.runs[0].text = translated_text
                    for r in paragraph.runs[1:]:
                        r.text = ""
                else:
                    paragraph.add_run().text = translated_text
    
    # Save the translated presentation
    print("保存翻译后的PPTX...")
    prs.save(output_path)
    
    # 输出详细的翻译统计
    print(f"✅ PPTX翻译完成！输出文件: {output_path}")
    print(f"📊 详细统计信息:")
    print(f"   - 总文本数: {len(text_items)}")
    print(f"   - 已翻译数: {translated_count}")
    print(f"   - 未翻译数: {untranslated_count}")
    print(f"   - 翻译完成率: {(translated_count/len(text_items)*100):.1f}%")
    print(f"   - 消耗 {total_token_count} tokens")
    print(f"   - 处理 {total_character_count} 字符")
    
    # Return metadata
    return {
        "translated_file_path": output_path,
        "token_count": total_token_count,
        "character_count": total_character_count,
        "total_texts": len(text_items),
        "translated_texts": translated_count,
        "untranslated_texts": untranslated_count,
        "translation_rate": (translated_count/len(text_items)*100)
    }

def main():
    if len(sys.argv) < 5:
        print("Usage: python translator_pptx_direct.py input.pptx output.pptx src_lang tgt_lang [engine] [max_workers]")
        sys.exit(1)
    inp = sys.argv[1]
    outp = sys.argv[2]
    src = sys.argv[3]
    tgt = sys.argv[4]
    engine = sys.argv[5] if len(sys.argv) > 5 else "deepseek"
    max_workers = int(sys.argv[6]) if len(sys.argv) > 6 else 4

    try:
        translate_pptx(inp, outp, src, tgt, engine=engine, max_workers=max_workers)
        print(f"✅ Translation finished. Output: {outp}")
    except Exception:
        print("❌ Translation failed:")
        traceback.print_exc()
        sys.exit(2)


if __name__ == "__main__":
    main()
